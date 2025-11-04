import io
from typing import Tuple, Dict

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Professional presentation styling
DEFAULT_FONTS = {
    'english': [
        'Arial Black',     # Extra bold, great for lyrics
        'Impact',         # Very bold, high visibility
        'Calibri Bold',   # Clean and bold
        'Segoe UI Black', # Modern bold
    ],
    'hindi': [
        'Noto Sans Devanagari Bold',
        'Mangal Bold',
        'Arial Unicode MS Bold',
    ],
    'telugu': [
        'Potti Sreeramulu',     # Primary choice
        'Noto Sans Telugu Bold',
        'Gautami Bold',
        'Vani Bold',
    ]
}

# Enhanced styling for maximum impact
DEFAULT_FONT_SIZE = Pt(48)  # Larger for better visibility
TEXT_COLOR = RGBColor(255, 255, 255)  # Bright white
BG_COLOR = RGBColor(0, 0, 0)  # Deep black for contrast
SHADOW_COLOR = RGBColor(0, 0, 0)  # Text shadow


def _detect_font(line: str, fonts: Dict[str, str]) -> str:
    """Pro font selection based on language and style preferences."""
    # Always prefer bold variants of fonts for lyrics
    
    if any(0x0900 <= ord(c) <= 0x097F for c in line):  # Hindi
        font = fonts.get('hindi', DEFAULT_FONTS['hindi'][0])
        # Try bold variant if not already specified
        if not any(bold in font.lower() for bold in ['bold', 'black']):
            font_bold = f"{font} Bold"
            return font_bold if font_bold in DEFAULT_FONTS['hindi'] else font
        return font
        
    if any(0x0C00 <= ord(c) <= 0x0C7F for c in line):  # Telugu
        font = fonts.get('telugu', DEFAULT_FONTS['telugu'][0])
        # For Telugu, maintain readability while being bold
        if not any(bold in font.lower() for bold in ['bold', 'black']):
            font_bold = f"{font} Bold"
            return font_bold if font_bold in DEFAULT_FONTS['telugu'] else font
        return font
    
    # English - prefer extra bold fonts for maximum impact
    font = fonts.get('english', DEFAULT_FONTS['english'][0])
    if not any(bold in font.lower() for bold in ['bold', 'black']):
        font_bold = f"{font} Bold"
        return font_bold if font_bold in DEFAULT_FONTS['english'] else font
    return font


def generate_pptx(
    lyrics_text: str,
    presentation_name: str,
    fonts: Dict[str, str] = None,
    font_size: int = None
) -> Tuple[bytes, int, str]:
    """
    Generate a professional PPTX file from lyrics text.

    Inputs:
    - lyrics_text: Full text with lines separated by newlines
    - presentation_name: Name without extension
    - fonts: Dict with 'english', 'hindi', 'telugu' font choices
    - font_size: Custom font size in points (default: 48)

    Returns:
    - pptx_bytes: Bytes content of the generated .pptx
    - slide_count: Number of slides generated
    - download_name: Suggested filename for download
    """
    fonts = fonts or {}
    font_size = Pt(font_size if font_size else DEFAULT_FONT_SIZE.pt)
    
    # Normalize and collect lines
    raw_lines = lyrics_text.splitlines()
    lines = [ln.strip() for ln in raw_lines if ln.strip()]  # Skip empty lines for better spacing

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]

    # Create slides with 2 lines each
    for i in range(0, len(lines), 2):
        slide = prs.slides.add_slide(blank_slide_layout)

        # Main lyrics textbox - positioned for maximum visibility
        left = Inches(0.5)
        top = Inches(2.0)  # Slightly higher
        width = Inches(12.33)  # Wider
        height = Inches(4.5)   # Taller
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        
        # Add lines with enhanced spacing
        current_lines = lines[i:i+2]
        if len(current_lines) == 2:
            # Two lines - add spacing between them
            tf.text = current_lines[0]
            p2 = tf.add_paragraph()
            p2.text = current_lines[1]
            p2.space_before = Pt(20)  # Space between lines
        else:
            # Single line - center vertically
            tf.text = current_lines[0]

        # Professional formatting for each paragraph
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = font_size
            p.font.bold = True  # Make text bold
            p.font.color.rgb = TEXT_COLOR
            p.font.name = _detect_font(p.text, fonts)

            # Add emphasis with slight line height adjustment
            p.line_spacing = 1.2  # 120% line spacing for better readability
            
            # Ensure consistent spacing
            p.space_after = 0
            if p.text:  # Only adjust spacing for non-empty paragraphs
                p.space_before = Pt(10)

        # Black background for maximum contrast
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    # Save to bytes
    bio = io.BytesIO()
    prs.save(bio)
    pptx_bytes = bio.getvalue()
    bio.close()

    slide_count = len(lines) // 2 + (1 if len(lines) % 2 else 0)
    download_name = f"{presentation_name}.pptx" if not presentation_name.lower().endswith('.pptx') else presentation_name

    return pptx_bytes, slide_count, download_name
